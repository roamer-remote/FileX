# Copyright (c) 2026 徐泽宇
"""pptx 业务逻辑模块。

Authors:
    徐泽宇
"""

from __future__ import annotations

try:
    from pptx import Presentation
except ModuleNotFoundError:
    Presentation = None

from services.extract.base import ExtractResult
from services.extract.loc_markers import format_slide_marker


def build_pptx_marked_body(path: str) -> str:
    if Presentation is None:
        raise ImportError("python-pptx is required to extract pptx files")
    prs = Presentation(path)
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    slide_lines.append(t)
        parts.append(format_slide_marker(slide_num))
        if slide_lines:
            parts.append("\n".join(slide_lines))
    return "\n\n".join(parts).strip()


def extract_pptx(path: str) -> ExtractResult:
    return ExtractResult(text=build_pptx_marked_body(path), engine="python-pptx")
